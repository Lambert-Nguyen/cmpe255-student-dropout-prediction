#!/usr/bin/env python3
"""
Build CMPE 255 Project Presentation (PPTX).

Generates a 19-slide, 16:9 widescreen presentation from project results.
All figures are pulled from results/figures/ and metrics from results/tables/.
"""

import os
import csv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent.parent
FIGURES = BASE_DIR / "results" / "figures"
TABLES = BASE_DIR / "results" / "tables"
OUTPUT = Path(__file__).resolve().parent / "CMPE255_Project.pptx"

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Colors ────────────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
SJSU_BLUE = RGBColor(0x00, 0x55, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
TEAL = RGBColor(0x1D, 0x9E, 0x75)
RED = RGBColor(0xE2, 0x4B, 0x4A)
AMBER = RGBColor(0xEF, 0x9F, 0x27)
ACCENT_BLUE = RGBColor(0x3B, 0x7D, 0xD8)

# ── Fonts ─────────────────────────────────────────────────────────────────────
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

# ── Metric data (extracted from actual notebook outputs) ──────────────────────
MODEL_DATA = [
    # Model, Accuracy, F1 Macro, F1 Weighted, CV F1 Mean, CV F1 Std
    ("XGBoost", 0.7695, 0.7056, 0.7634, 0.7136, 0.0131),
    ("Random Forest", 0.7706, 0.6887, 0.7548, 0.6956, 0.0087),
    ("Logistic Regression", 0.7684, 0.6826, 0.7531, 0.6782, 0.0089),
    ("SVM", 0.7593, 0.6797, 0.7466, 0.6798, 0.0220),
    ("Decision Tree", 0.7119, 0.6346, 0.7052, 0.6443, 0.0101),
    ("KNN (k=11)", 0.6949, 0.5860, 0.6694, 0.6036, 0.0094),
]

ASSOC_RULES = [
    ("CU1_Low, Age_Mid", "Target_Dropout", 0.058, 0.874, 2.72),
    ("AdmGrade_Low, CU1_Low", "Target_Dropout", 0.069, 0.860, 2.68),
    ("CU1_Low", "Target_Dropout", 0.175, 0.771, 2.40),
    ("CU1_Low, AdmGrade_Med", "Target_Dropout", 0.084, 0.764, 2.38),
    ("TuitionUpToDate, CU1_Low", "Target_Dropout", 0.115, 0.696, 2.17),
    ("TuitionUpToDate, CU1_High, Scholarship, AdmGrade_Med", "Age_Young, Target_Graduate", 0.078, 0.801, 2.09),
]


# ══════════════════════════════════════════════════════════════════════════════
#  Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_header_bar(slide, title_text):
    """Add a dark-blue header bar with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)


def add_slide_number(slide, num, total=19):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45), Inches(1.0), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


def add_body_text(slide, left, top, width, height, bullets, font_size=18,
                  bold_first_word=False, line_spacing=1.3):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
        p.level = 0
        if bold_first_word and ":" in text:
            # We can't easily bold partial runs after setting p.text,
            # so we rebuild with runs
            p.clear()
            parts = text.split(":", 1)
            run1 = p.add_run()
            run1.text = parts[0] + ":"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = DARK_GRAY
            run1.font.name = FONT_BODY
            if len(parts) > 1:
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = DARK_GRAY
                run2.font.name = FONT_BODY
    return txBox


def add_image(slide, img_name, left, top, width=None, height=None):
    """Insert an image from the figures directory."""
    img_path = str(FIGURES / img_name)
    if not os.path.exists(img_path):
        # Add placeholder text instead
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(4), Inches(1)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image not found: {img_name}]"
        p.font.size = Pt(14)
        p.font.color.rgb = RED
        return None

    kwargs = {}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(img_path, Inches(left), Inches(top), **kwargs)


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_accent_box(slide, left, top, width, height, text, font_size=16):
    """Add a teal accent callout box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = TEAL
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, width, height, header, rows, col_widths=None):
    """Add a formatted table to the slide."""
    n_rows = len(rows) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                         Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_dim_box(slide, left, top, width, height, title, items, color=SJSU_BLUE):
    """Add a dimension box for the star schema."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)

    return box


def add_connector_line(slide, x1, y1, x2, y2, color=SJSU_BLUE):
    """Add a simple line connector."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


# ══════════════════════════════════════════════════════════════════════════════
#  Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def build_slide_01_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full-slide dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Teal accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.35), Inches(10.3), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Predicting Student Dropout and\nAcademic Success Using Data Mining"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "CMPE 255 — Data Mining  |  Spring 2026"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEAL
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.LEFT

    # Team info
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    lines = [
        ("Lam Nguyen (018229432)  •  Tri Ngo", Pt(18)),
        ("San José State University", Pt(16)),
        ("April 2026", Pt(16)),
    ]
    for i, (text, size) in enumerate(lines):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)

    add_slide_number(slide, 1)
    add_speaker_notes(slide, (
        "Welcome everyone. Today we present our CMPE 255 project on predicting "
        "student dropout using data mining techniques. We analyzed a dataset of "
        "over 4,400 students from a Portuguese university to identify key risk factors "
        "and build predictive models."
    ))


def build_slide_02_intro(prs):
    """Slide 2: Introduction / Motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Introduction & Motivation")

    bullets = [
        "Higher education dropout is a significant global challenge — "
        "it affects students, families, and institutions alike",
        "~32% of students in our dataset dropped out before graduating",
        "Early identification of at-risk students enables targeted academic "
        "advising and intervention programs",
        "Data mining can uncover hidden patterns in enrollment, demographic, "
        "and academic performance data",
    ]
    add_body_text(slide, 0.6, 1.5, 7.0, 3.5, bullets, font_size=19)

    add_accent_box(
        slide, 0.6, 5.3, 12.1, 0.9,
        'Goal: "Apply classification, clustering, and association rule mining '
        'to predict student outcomes and identify key dropout risk factors"',
        font_size=17
    )

    # Right side — class distribution chart thumbnail
    add_image(slide, "class_distribution.png", 8.2, 1.4, width=4.5)

    add_slide_number(slide, 2)
    add_speaker_notes(slide, (
        "Student dropout is costly — both for students who invest time and money, "
        "and for universities that lose funding and reputation. In our dataset, "
        "nearly one-third of students dropped out. Our goal is to use data mining "
        "to build an early-warning system that flags at-risk students so advisors "
        "can intervene before it's too late."
    ))


def build_slide_03_problem(prs):
    """Slide 3: Problem Statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Problem Statement")

    bullets = [
        "Research Question: Can we predict whether a student will drop out, "
        "remain enrolled, or graduate based on enrollment-time data?",
        "Target Variable: 3-class classification — Dropout / Enrolled / Graduate",
        "Scope: Supervised classification + unsupervised clustering for risk profiling",
        "Challenge: Class imbalance — Graduate ~50%, Dropout ~32%, Enrolled ~18%",
        "Expected Outcome: An early-warning framework to flag at-risk students "
        "for academic advising",
    ]
    add_body_text(slide, 0.6, 1.5, 12.0, 4.5, bullets, font_size=19,
                  bold_first_word=True)

    add_slide_number(slide, 3)
    add_speaker_notes(slide, (
        "Our core research question is whether we can predict student outcomes "
        "using data available at enrollment time and after the first semesters. "
        "This is a multi-class problem with three outcomes. The class imbalance "
        "is a key challenge — Enrolled students make up only 18% of the data, "
        "making them harder to predict correctly."
    ))


def build_slide_04_data(prs):
    """Slide 4: Data Sources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Data Sources")

    # Main info
    info_bullets = [
        "Source: UCI Machine Learning Repository (Dataset ID: 697)",
        "Origin: Instituto Politécnico de Portalegre, Portugal",
        "Size: 4,424 students × 37 columns (36 features + 1 target)",
        "Missing values: None  |  License: CC BY 4.0",
        "Citation: Realinho et al. (2021)",
    ]
    add_body_text(slide, 0.6, 1.4, 5.5, 2.5, info_bullets, font_size=16)

    # Feature category boxes
    categories = [
        ("Demographic", ["Marital status", "Nationality", "Gender", "Age at enrollment"], SJSU_BLUE),
        ("Academic", ["Previous qualification", "Admission grade", "Course", "Application mode"], TEAL),
        ("Socio-Economic", ["Scholarship holder", "Debtor", "Tuition fees", "Parental education"], AMBER),
        ("Performance", ["Units enrolled (sem 1 & 2)", "Units approved (sem 1 & 2)", "Grades (sem 1 & 2)"], RED),
    ]

    x_positions = [0.6, 3.5, 6.4, 9.3]
    for i, (title, items, color) in enumerate(categories):
        add_dim_box(slide, x_positions[i], 4.3, 2.7, 2.5, title, items, color)

    add_slide_number(slide, 4)
    add_speaker_notes(slide, (
        "Our dataset comes from the UCI repository — compiled from institutional "
        "databases at a Portuguese polytechnic. It contains 4,424 student records "
        "with 36 features spanning demographics, prior academics, socio-economic "
        "factors, and semester performance. Importantly, there are no missing values, "
        "so we could focus on modeling rather than imputation."
    ))


def build_slide_05_star_schema(prs):
    """Slide 5: Data Warehouse Design (Star Schema)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Data Warehouse Design — Star Schema")

    # Fact table (center)
    fact_box = add_dim_box(
        slide, 4.5, 2.8, 4.3, 2.5,
        "Student_Enrollment_Fact",
        ["student_id (PK)", "course_id (FK)", "sem1_grade, sem2_grade",
         "sem1_approved, sem2_approved", "target_outcome"],
        DARK_BLUE
    )

    # Dimension tables around the fact
    dims = [
        (0.4, 1.5, 3.3, 2.0, "Student_Dim",
         ["age", "gender", "marital_status", "nationality"]),
        (9.6, 1.5, 3.3, 2.0, "Academic_Dim",
         ["prev_qualification", "application_mode", "admission_grade", "course"]),
        (0.4, 4.8, 3.3, 2.0, "Family_Dim",
         ["mother_qualification", "father_qualification", "mother_occupation", "father_occupation"]),
        (9.6, 4.8, 3.3, 2.0, "Economic_Dim",
         ["scholarship", "debtor", "tuition_fees", "unemployment_rate, GDP"]),
    ]
    for (x, y, w, h, title, items) in dims:
        add_dim_box(slide, x, y, w, h, title, items, SJSU_BLUE)

    # Connectors from dims to fact
    add_connector_line(slide, 3.7, 2.5, 4.5, 3.5)   # Student -> Fact
    add_connector_line(slide, 9.6, 2.5, 8.8, 3.5)   # Academic -> Fact
    add_connector_line(slide, 3.7, 5.0, 4.5, 4.5)   # Family -> Fact
    add_connector_line(slide, 9.6, 5.8, 8.8, 4.5)   # Economic -> Fact

    # ETL note
    txBox = slide.shapes.add_textbox(Inches(3.5), Inches(1.3), Inches(6.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ETL: Source DBs → Extract → Transform (encode, normalize) → Load → Star Schema"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 5)
    add_speaker_notes(slide, (
        "Here we show a conceptual star schema for our data warehouse design. "
        "The central fact table stores each student's enrollment record with "
        "semester grades and outcome. Four dimension tables capture student "
        "demographics, academic background, family information, and economic factors. "
        "The ETL pipeline encodes categorical features and normalizes numerics."
    ))


def build_slide_06_preprocessing(prs):
    """Slide 6: Data Preprocessing Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Data Preprocessing Pipeline")

    steps = [
        ("1. Inspection", "4,424 rows, 36 features, 0 missing values"),
        ("2. Feature Typing", "Identified categorical (int-coded) vs continuous"),
        ("3. Target Encoding", "Dropout → 0, Enrolled → 1, Graduate → 2"),
        ("4. Normalization", "StandardScaler on all numeric features"),
        ("5. Stratified Split", "80% train (3,539) / 20% test (885)"),
    ]

    # Draw pipeline as connected boxes
    box_w, box_h = 2.2, 1.2
    start_x = 0.5
    y = 2.8
    for i, (step, desc) in enumerate(steps):
        x = start_x + i * 2.5

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        color = TEAL if i < 4 else SJSU_BLUE
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        p2.font.name = FONT_BODY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

        # Arrow between boxes
        if i < len(steps) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(y + 0.4),
                Inches(0.3), Inches(0.35)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # Additional notes below
    notes = [
        "random_state=42 for full reproducibility across all operations",
        "Stratified split preserves class ratios in both train and test sets",
        "Unscaled copy saved separately for EDA visualizations",
    ]
    add_body_text(slide, 0.6, 4.8, 12.0, 2.0, notes, font_size=15)

    add_slide_number(slide, 6)
    add_speaker_notes(slide, (
        "Our preprocessing pipeline has five steps. First we verified data quality — "
        "no missing values. We identified feature types, encoded the target variable "
        "numerically, standardized all features, and performed a stratified 80/20 split "
        "to preserve the class distribution in both train and test sets."
    ))


def build_slide_07_eda_distribution(prs):
    """Slide 7: Class Distribution & Descriptive Stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "EDA — Class Distribution")

    add_image(slide, "class_distribution.png", 0.4, 1.4, width=6.0)

    stats_bullets = [
        "Graduate: 2,209 students (49.9%)",
        "Dropout: 1,421 students (32.1%)",
        "Enrolled: 794 students (17.9%)",
        "Total: 4,424 students",
        "Notable imbalance — Enrolled class is underrepresented",
    ]
    add_body_text(slide, 7.0, 1.8, 5.5, 3.5, stats_bullets, font_size=17)

    add_accent_box(
        slide, 7.0, 5.2, 5.5, 0.8,
        "Stratified sampling ensures all classes represented in train/test",
        font_size=14
    )

    add_slide_number(slide, 7)
    add_speaker_notes(slide, (
        "Our target distribution shows moderate class imbalance. About half the "
        "students graduated, a third dropped out, and only 18% remain enrolled. "
        "This imbalance affects model performance, particularly for the Enrolled "
        "class. We use stratified splitting and F1-macro scoring to handle this."
    ))


def build_slide_08_eda_features(prs):
    """Slide 8: Feature Distributions & Correlations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "EDA — Feature Distributions & Correlations")

    add_image(slide, "feature_distributions.png", 0.3, 1.3, width=6.5)
    add_image(slide, "top_correlated_features.png", 6.9, 1.3, width=6.0)

    add_accent_box(
        slide, 0.6, 6.0, 12.1, 0.7,
        "Key Insight: Students who drop out have dramatically fewer approved "
        "curricular units in both semesters",
        font_size=15
    )

    add_slide_number(slide, 8)
    add_speaker_notes(slide, (
        "The KDE plots on the left show clear separation — dropout students cluster "
        "near zero approved units while graduates spread across higher values. "
        "The Spearman correlation chart on the right confirms that 2nd semester "
        "approved units is the single strongest predictor of student outcome."
    ))


def build_slide_09_eda_deeper(prs):
    """Slide 9: Deeper EDA Insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "EDA — Deeper Insights")

    add_image(slide, "boxplots_grades.png", 0.2, 1.3, width=4.3)
    add_image(slide, "stacked_bar_binary.png", 4.6, 1.3, width=4.3)
    add_image(slide, "pca_scatter.png", 9.0, 1.3, width=4.1)

    captions = [
        ("Semester grades by outcome", 0.5, 5.5),
        ("Financial factors vs outcome", 5.0, 5.5),
        ("PCA 2D projection", 9.5, 5.5),
    ]
    for text, x, y in captions:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

    add_accent_box(
        slide, 0.6, 6.1, 12.1, 0.7,
        "PCA shows moderate class separability; scholarship holders show "
        "significantly lower dropout rates",
        font_size=15
    )

    add_slide_number(slide, 9)
    add_speaker_notes(slide, (
        "Three key views here. Left: boxplots show dramatic separation in semester "
        "grades — dropout students score much lower. Center: stacked bars reveal "
        "that being a debtor or not having tuition up to date increases dropout risk. "
        "Right: PCA projects all 36 features to 2D, showing clusters but with overlap, "
        "suggesting a non-trivial classification problem."
    ))


def build_slide_10_viz_summary(prs):
    """Slide 10: Visualization Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Visualization Techniques Summary")

    # 2×3 grid of thumbnail charts
    images = [
        ("correlation_heatmap.png", "Correlation Heatmap"),
        ("feature_distributions.png", "KDE Distributions"),
        ("boxplots_grades.png", "Boxplots"),
        ("stacked_bar_binary.png", "Stacked Bar Charts"),
        ("pca_scatter.png", "PCA Scatter Plot"),
        ("rf_feature_importance.png", "Feature Importance"),
    ]

    positions = [
        (0.4, 1.4), (4.5, 1.4), (8.6, 1.4),
        (0.4, 3.8), (4.5, 3.8), (8.6, 3.8),
    ]

    for (img_name, label), (x, y) in zip(images, positions):
        add_image(slide, img_name, x, y, width=3.8, height=2.1)
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 2.1), Inches(3.8), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

    # Caption at bottom
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ("All visualizations generated with Matplotlib & Seaborn  |  "
              "Consistent color coding: Red = Dropout, Amber = Enrolled, Green = Graduate")
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 10)
    add_speaker_notes(slide, (
        "This slide showcases the breadth of visualization techniques we used — "
        "heatmaps, KDE distributions, boxplots, stacked bars, PCA, and feature "
        "importance charts. All use a consistent color scheme throughout the project."
    ))


def build_slide_11_classification_overview(prs):
    """Slide 11: Classification Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Classification Models — Overview")

    models = [
        ("Decision Tree", "max_depth=10"),
        ("Random Forest", "n_estimators=200"),
        ("KNN", "best k=11 (tuned via accuracy)"),
        ("Logistic Regression", "multinomial, max_iter=1000"),
        ("SVM", "RBF kernel, probability=True"),
        ("XGBoost", "n_estimators=200, max_depth=6, lr=0.1"),
    ]

    # Model boxes in 2x3 grid
    for i, (name, params) in enumerate(models):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.5 + row * 2.2

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(3.8), Inches(1.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = SJSU_BLUE
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = SJSU_BLUE
        p.font.name = FONT_BODY

        p2 = tf.add_paragraph()
        p2.text = params
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = FONT_BODY
        p2.space_before = Pt(6)

    # Bottom note
    add_accent_box(
        slide, 0.6, 6.0, 12.1, 0.7,
        "All models: random_state=42  |  80/20 stratified split  |  "
        "5-fold cross-validation  |  F1-macro as primary metric",
        font_size=14
    )

    add_slide_number(slide, 11)
    add_speaker_notes(slide, (
        "We trained six classification models, covering tree-based, instance-based, "
        "linear, kernel, and boosting approaches. All use random_state=42 for "
        "reproducibility. KNN's k was tuned by testing 3, 5, 7, 9, and 11, with "
        "k=11 yielding the best accuracy. We evaluate all models using F1-macro "
        "to account for class imbalance."
    ))


def build_slide_12_classification_results(prs):
    """Slide 12: Classification Results (Confusion Matrices & ROC)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Classification Results — Best vs Worst")

    # Best model (XGBoost) CM and ROC
    add_image(slide, "xgb_cm.png", 0.3, 1.3, width=3.7, height=2.8)
    add_image(slide, "xgb_roc.png", 0.3, 4.2, width=3.7, height=2.8)

    # Mid model (Logistic Regression) CM
    add_image(slide, "lr_cm.png", 4.3, 1.3, width=3.7, height=2.8)
    add_image(slide, "lr_roc.png", 4.3, 4.2, width=3.7, height=2.8)

    # Worst model (KNN) CM
    add_image(slide, "knn_cm.png", 8.3, 1.3, width=3.7, height=2.8)
    add_image(slide, "knn_roc.png", 8.3, 4.2, width=3.7, height=2.8)

    add_slide_number(slide, 12)
    add_speaker_notes(slide, (
        "Comparing our best model (XGBoost, F1-macro 0.706), a mid-range model "
        "(Logistic Regression, 0.683), and the weakest (KNN, 0.586). The confusion "
        "matrices show that all models struggle most with the Enrolled class due to "
        "its small size. ROC curves show strong AUC values across all three classes, "
        "especially for Dropout and Graduate."
    ))


def build_slide_13_clustering(prs):
    """Slide 13: Clustering."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Clustering Analysis")

    add_image(slide, "kmeans_elbow_silhouette.png", 0.3, 1.3, width=6.0, height=2.6)
    add_image(slide, "kmeans_crosstab.png", 0.3, 4.1, width=5.5, height=2.8)
    add_image(slide, "hierarchical_dendrogram.png", 6.5, 1.3, width=6.3, height=2.6)

    # Clustering insights
    insights = [
        "K-Means (k=3): Silhouette = 0.2145",
        "Clusters partially recover outcome classes",
        "Hierarchical (Ward): Silhouette = 0.1158",
        "Low scores confirm hard boundary problem",
        "Supervision needed for reliable prediction",
    ]
    add_body_text(slide, 6.8, 4.2, 5.5, 2.8, insights, font_size=15)

    add_slide_number(slide, 13)
    add_speaker_notes(slide, (
        "K-Means with k=3 shows the best silhouette score at 0.21, confirming "
        "three natural groupings. However, the crosstab shows significant overlap — "
        "clusters don't cleanly separate the three outcome classes. Hierarchical "
        "clustering with Ward linkage confirms this. The low silhouette scores "
        "indicate that unsupervised methods alone cannot reliably predict outcomes, "
        "justifying our supervised classification approach."
    ))


def build_slide_14_association(prs):
    """Slide 14: Association Rules."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Association Rule Mining (Apriori)")

    # Explanation
    bullets = [
        "Binned numeric features into categorical ranges",
        "Applied Apriori: min_support=0.05, min_confidence=0.60",
        "189 frequent itemsets → 334 rules → 85 with Target as consequent",
    ]
    add_body_text(slide, 0.6, 1.4, 12.0, 1.5, bullets, font_size=16)

    # Top rules table
    header = ["Antecedent", "Consequent", "Support", "Conf.", "Lift"]
    rows = []
    for ante, cons, sup, conf, lift in ASSOC_RULES:
        rows.append([ante, cons, f"{sup:.3f}", f"{conf:.3f}", f"{lift:.2f}"])

    add_table(slide, 0.5, 3.2, 12.3, 3.0, header, rows,
              col_widths=[4.0, 3.0, 1.3, 1.3, 1.3])

    add_accent_box(
        slide, 0.6, 6.4, 12.1, 0.6,
        "Top finding: Low 1st-semester approved units + Mid age → Dropout "
        "(confidence = 87.4%, lift = 2.72)",
        font_size=14
    )

    add_slide_number(slide, 14)
    add_speaker_notes(slide, (
        "We applied the Apriori algorithm after binning numeric features. "
        "The strongest association rule says: students with low first-semester "
        "approved units who are mid-age have an 87% probability of dropping out, "
        "with a lift of 2.72 — meaning they are 2.7 times more likely to drop out "
        "than the baseline. This is actionable intelligence for academic advisors."
    ))


def build_slide_15_comparison(prs):
    """Slide 15: Model Comparison Table (MOST IMPORTANT SLIDE)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Model Comparison — Results Summary")

    # Table with all model results
    header = ["Model", "Accuracy", "F1 (Macro)", "F1 (Weighted)", "CV F1 Mean ± Std"]
    rows = []
    for name, acc, f1m, f1w, cvm, cvs in MODEL_DATA:
        rows.append([
            name,
            f"{acc:.1%}",
            f"{f1m:.4f}",
            f"{f1w:.4f}",
            f"{cvm:.4f} ± {cvs:.4f}",
        ])

    table_shape = add_table(slide, 0.5, 1.4, 7.5, 3.8, header, rows,
                            col_widths=[2.0, 1.2, 1.2, 1.2, 1.9])

    # Highlight best values (manually mark the best model row)
    table = table_shape.table
    # XGBoost row (row 1) — best F1-macro
    for j in range(len(header)):
        cell = table.cell(1, j)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = TEAL

    # Random Forest row (row 2) — best accuracy
    cell_acc = table.cell(2, 1)
    p_acc = cell_acc.text_frame.paragraphs[0]
    p_acc.font.bold = True
    p_acc.font.color.rgb = TEAL

    # Insert model comparison chart
    add_image(slide, "model_comparison.png", 8.3, 1.4, width=4.6, height=3.5)

    # Key takeaway
    add_accent_box(
        slide, 0.6, 5.8, 12.1, 1.0,
        "XGBoost achieves the highest F1-macro (0.7056) with strong "
        "cross-validation (0.7136 ± 0.0131)\n"
        "Random Forest leads in accuracy (77.1%) with the most stable CV (± 0.0087)",
        font_size=14
    )

    add_slide_number(slide, 15)
    add_speaker_notes(slide, (
        "This is our key results slide. XGBoost achieves the best F1-macro score "
        "of 0.706, which is our primary metric since it accounts for class imbalance. "
        "Random Forest has the highest raw accuracy at 77.1% and the most stable "
        "cross-validation performance. KNN and Decision Tree trail behind. "
        "The cross-validation results confirm these models generalize well — "
        "the small standard deviations show consistent performance across folds."
    ))


def build_slide_16_feature_importance(prs):
    """Slide 16: Feature Importance & Key Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Feature Importance & Key Findings")

    add_image(slide, "feature_importance.png", 0.2, 1.3, width=8.0, height=4.5)

    # Key features callout
    features = [
        "#1: Curricular units 2nd sem (approved)",
        "#2: Curricular units 2nd sem (grade)",
        "#3: Curricular units 1st sem (approved)",
        "#4: Tuition fees up to date",
        "#5: Age at enrollment",
    ]
    txBox = add_body_text(slide, 8.5, 1.5, 4.3, 3.5, features, font_size=15,
                          bold_first_word=True)

    add_accent_box(
        slide, 8.5, 4.8, 4.3, 1.8,
        "Both Random Forest and XGBoost\n"
        "agree: semester performance is\n"
        "the dominant predictor of\n"
        "student outcomes",
        font_size=14
    )

    add_slide_number(slide, 16)
    add_speaker_notes(slide, (
        "Feature importance analysis from both Random Forest and XGBoost consistently "
        "ranks 2nd semester approved units as the single strongest predictor. "
        "Academic performance features dominate the top 5, followed by financial "
        "indicators like tuition status. This tells us that early semester performance "
        "is the most reliable signal for identifying at-risk students."
    ))


def build_slide_17_conclusions(prs):
    """Slide 17: Conclusions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Conclusions")

    findings = [
        "XGBoost achieved the highest F1-macro of 0.7056, outperforming "
        "all other models including Random Forest (0.6887) and Logistic Regression (0.6826)",
        "1st and 2nd semester academic performance (approved units, grades) "
        "are the dominant predictors of student outcome",
        "Association rules reveal: students with ≤2 approved units in semester 1 "
        "have a 77% dropout probability (lift = 2.4×)",
        "K-Means clustering reveals natural risk profiles that partially align "
        "with outcome labels (silhouette = 0.21)",
    ]
    add_body_text(slide, 0.6, 1.5, 12.0, 3.5, findings, font_size=17)

    # Recommendation box
    add_accent_box(
        slide, 1.0, 5.0, 11.3, 1.0,
        "Recommendation: Universities should flag students with fewer than 3 "
        "approved units after semester 1 for immediate academic advising intervention",
        font_size=17
    )

    # Summary line
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ('"Data mining techniques can reliably identify at-risk students '
              'early enough for meaningful intervention."')
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = SJSU_BLUE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 17)
    add_speaker_notes(slide, (
        "To summarize: XGBoost is our best model by F1-macro, semester performance "
        "is the strongest predictor, and association rules provide actionable "
        "thresholds for intervention. Our key recommendation is simple — track "
        "first-semester approved units and flag students with fewer than 3 for "
        "immediate advising."
    ))


def build_slide_18_future(prs):
    """Slide 18: Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Future Work")

    items = [
        "Real-time prediction: Incorporate live semester data for dynamic "
        "risk scoring throughout the academic year",
        "Deployment: Build a web dashboard (Streamlit/Flask) for academic "
        "advisors to query student risk in real time",
        "Multi-institution validation: Test model generalizability across "
        "universities in different countries",
        "Deep learning: Explore neural network architectures (e.g., TabNet) "
        "for potential performance gains",
        "Temporal modeling: Apply sequence models to track how student risk "
        "evolves semester by semester",
    ]
    add_body_text(slide, 0.6, 1.5, 12.0, 5.0, items, font_size=18,
                  bold_first_word=True)

    # Framing note
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Limitations → Opportunities: Each future direction addresses a current constraint"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_BODY

    add_slide_number(slide, 18)
    add_speaker_notes(slide, (
        "Looking ahead, the most impactful next step would be building a real-time "
        "dashboard that advisors can use during the semester. We'd also like to "
        "validate our models across multiple institutions. Deep learning approaches "
        "like TabNet could potentially improve performance, and temporal modeling "
        "could track how risk evolves over time."
    ))


def build_slide_19_thankyou(prs):
    """Slide 19: Thank You / Questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Teal accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.3), Inches(6.3), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Thank you
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER

    # Questions
    txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.3), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Questions?"
    p2.font.size = Pt(32)
    p2.font.color.rgb = TEAL
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.CENTER

    # Contact info
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(4.8), Inches(9.3), Inches(2.0))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    lines = [
        "Lam Nguyen (018229432)  •  Tri Ngo",
        "CMPE 255 — Data Mining  |  Spring 2026",
        "San José State University",
        "",
        "github.com/duylam1407/CMPE255-Student-Dropout-Prediction",
    ]
    for i, text in enumerate(lines):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if i < 3 else Pt(14)
        p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)

    add_slide_number(slide, 19)
    add_speaker_notes(slide, (
        "Thank you for your attention. We're happy to take any questions about "
        "our methodology, results, or potential applications of this work."
    ))


# ══════════════════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 19 slides
    build_slide_01_title(prs)       # Title
    build_slide_02_intro(prs)       # Introduction / Motivation
    build_slide_03_problem(prs)     # Problem Statement
    build_slide_04_data(prs)        # Data Sources
    build_slide_05_star_schema(prs) # Data Warehouse Design
    build_slide_06_preprocessing(prs)  # Data Preprocessing
    build_slide_07_eda_distribution(prs)  # EDA — Class Distribution
    build_slide_08_eda_features(prs)      # EDA — Features & Correlations
    build_slide_09_eda_deeper(prs)        # EDA — Deeper Insights
    build_slide_10_viz_summary(prs)       # Visualization Summary
    build_slide_11_classification_overview(prs)  # Classification Overview
    build_slide_12_classification_results(prs)   # Classification Results
    build_slide_13_clustering(prs)        # Clustering
    build_slide_14_association(prs)       # Association Rules
    build_slide_15_comparison(prs)        # Model Comparison (KEY SLIDE)
    build_slide_16_feature_importance(prs)  # Feature Importance
    build_slide_17_conclusions(prs)       # Conclusions
    build_slide_18_future(prs)            # Future Work
    build_slide_19_thankyou(prs)          # Thank You

    prs.save(str(OUTPUT))
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
